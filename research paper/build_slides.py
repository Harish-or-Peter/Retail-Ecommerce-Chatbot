"""Generate presentation.pptx (13 slides, figures embedded) for the retail e-commerce paper.

Run from this folder:
    pip install python-pptx
    python build_slides.py
Produces presentation.pptx next to this script.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
FIG  = os.path.join(BASE, "resources", "figures")
DIA  = os.path.join(BASE, "resources", "diagrams")
OUT  = os.path.join(BASE, "presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
NAVY = RGBColor(0x10, 0x3a, 0x5a); DARK = RGBColor(0x22, 0x22, 0x22)

def slide(title):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    bar.text_frame.margin_left = Inches(0.4)
    p = bar.text_frame.paragraphs[0]; p.text = title
    p.runs[0].font.size = Pt(28); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RGBColor(0xff,0xff,0xff)
    return s

def bullets(s, items, left=0.6, top=1.4, width=12.0, size=18):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.6))
    tf = box.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = it[1] if isinstance(it, tuple) else 0
        txt = it[0] if isinstance(it, tuple) else it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "   – ") + txt; p.level = lvl
        p.runs[0].font.size = Pt(size if lvl == 0 else size-2); p.runs[0].font.color.rgb = DARK
        p.space_after = Pt(8)

def pic(s, path, left, top, width):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))

# 1 Title
s = prs.slides.add_slide(BLANK)
box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3)); tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = ("Fine-Tuning a Compact Instruction-Tuned Language Model for "
                                "Retail E-Commerce Customer-Support Conversational AI")
p.runs[0].font.size = Pt(34); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
for t in ["Harish Chandra",
          "Master's in CS: Artificial Intelligence and Machine Learning, Woolf University", "31 May 2026"]:
    q = tf.add_paragraph(); q.text = t; q.runs[0].font.size = Pt(18); q.runs[0].font.color.rgb = DARK; q.space_before = Pt(10)

bullets(slide("Problem & Research Question"), [
    "E-commerce support is high-volume and repetitive: orders, refunds, cancellations, returns, delivery, payments.",
    "Customers expect instant, 24/7 answers - costly to staff with humans alone.",
    "Large LLMs work but are expensive; out of reach for many small/mid-sized retailers.",
    "Research question:",
    ("Can flan-t5-small be fine-tuned on retail support data to give accurate, relevant answers on a single free-tier GPU?", 1),
])

bullets(slide("Why Retail E-Commerce?"), [
    "Online retail is a multi-trillion-dollar, fast-growing sector; support volume scales with sales.",
    "Chatbot market ~$7.76B (2024) -> ~$27.29B (2030); customer service is the largest application.",
    "A small set of recurring intents dominates the workload - highly automatable.",
    "Business value: lower cost-per-contact, instant 24/7 responses, consistency, scalability.",
])

bullets(slide("Literature & Research Gap"), [
    "Foundations: Transformer (Vaswani 2017), T5 (Raffel 2020), instruction tuning / FLAN (Chung 2022).",
    "Applied: fine-tuned transformers & generative assistants for support (Nandi 2024; Khennouche 2023).",
    "Gap 1: low-resource, compact models for retail support are under-reported.",
    "Gap 2: most work isolates intent classification, not end-to-end response generation.",
    "Contribution: reproducible single-GPU end-to-end study with honest failure analysis.",
])

s = slide("Methodology - Pipeline")
bullets(s, ["Fine-tuning only (no RAG).", "flan-t5-small (~77M params).",
            "FP32 on Colab T4 (FP16 NaNs; no bf16).", "5 epochs, early stopping, seed 42.",
            "Beam search (4 beams) at inference."], left=0.6, top=1.4, width=6.0)
pic(s, os.path.join(DIA, "pipeline_flow.png"), 7.0, 1.3, 5.6)

s = slide("Dataset & EDA")
bullets(s, ["Bitext retail e-commerce (HF, CDLA, synthetic - no PII).", "44,884 instruction-response pairs.",
            "13 categories, 46 intents.", "Split 35,907 / 4,488 / 4,489 (stratified).",
            "Token lengths -> input 128 / target 256."], left=0.6, top=1.4, width=6.2)
pic(s, os.path.join(FIG, "fig2_intent_distribution.png"), 7.0, 1.5, 5.8)

bullets(slide("Fine-Tuning Setup"), [
    "AdamW, learning rate 3e-4, weight decay 0.01, 500-step linear warmup.",
    "Batch 8 x grad-accum 2 = effective 16.",
    "Labels padded with -100 so loss ignores padding.",
    "Seq2SeqTrainer; best checkpoint by validation loss.",
    "Up to 5 epochs - within the 25-epoch budget.",
])

s = slide("Results - Baseline vs Fine-Tuned")
bullets(s, ["ROUGE-L: 0.057 -> 0.453 (~8x).", "BLEU: 0.0 -> 29.27.",
            "BERTScore-F1: 0.818 -> 0.909.", "Avg response: 7 -> 133 words."], left=0.6, top=1.4, width=6.0)
pic(s, os.path.join(FIG, "fig5_metric_comparison.png"), 7.0, 1.5, 5.8)

s = slide("Results - Convergence & Per-Category")
pic(s, os.path.join(FIG, "fig4_training_loss.png"), 0.5, 1.5, 6.2)
pic(s, os.path.join(FIG, "fig6_per_category_rougeL.png"), 6.9, 1.5, 6.0)

s = slide("Demonstration & Qualitative Behaviour")
bullets(s, [
    "Gradio chat interface served from the notebook.",
    "In-domain: structured, step-by-step answers.",
    "Baseline gave useless one-liners.",
    "Ambiguous words (e.g. 'refund') -> inferred intent.",
    "Out-of-domain -> answered in support style (over-generalisation).",
], left=0.6, top=1.4, width=6.2, size=17)
pic(s, os.path.join(DIA, "gradio_demo4.png"), 6.9, 2.2, 6.1)

bullets(slide("Discussion & Limitations"), [
    "Domain adaptation - not scale - was the limiting factor.",
    "Practical: cheap to reproduce/retrain; automate high-volume intents first.",
    "Limitation: no out-of-scope detection -> needs intent filter + human fallback.",
    "Limitation: synthetic data; automatic metrics approximate human judgement.",
])

bullets(slide("Conclusion & Future Work"), [
    "Compact, free-tier-trainable model handles routine retail support well (8x ROUGE-L gain).",
    "Contribution: transparent, reproducible blueprint for low-resource domain chatbots.",
    "Future: out-of-scope detection + human handoff; real logs + human eval;",
    "larger-model comparison (flan-t5-base); multi-turn; multilingual.",
])

bullets(slide("Key References"), [
    "Vaswani et al. (2017) - Attention Is All You Need.",
    "Raffel et al. (2020) - T5.",
    "Chung et al. (2022) - FLAN.",
    "Nandi et al. (2024); Khennouche et al. (2023) - customer-service chatbots.",
    "Bitext Retail E-Commerce dataset (HF, CDLA-Sharing-1.0).",
], size=16)

prs.save(OUT)
print("WROTE", OUT, "with", len(prs.slides._sldIdLst), "slides")
